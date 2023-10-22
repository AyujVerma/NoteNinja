from pptx import Presentation
import KeywordExtract
import re

delimiter_pattern = r'[.!?\n]'

def extract_notes_from_slides(file_path):
    all_sents = ""
    presentation = Presentation(file_path)
    slide_dict = dict()
    title_dict = dict()
    body_text = []
    images = []
    for slide in presentation.slides:
        # ignores title slide
        if (slide.slide_layout.name != 'Title Slide'):
            frame_counter = 0
            no_title_count = 1
            list_key = list()
            title = "No Title " + str(no_title_count)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.has_text_frame:
                        if frame_counter == 0:
                            title = shape.text
                        else:
                            sent_list = re.split(delimiter_pattern, shape.text)
                            for sent in sent_list:
                                if (len(sent) > 0 and sent != None):
                                    text_pair = [None, sent]
                                    body_text.append(text_pair)
                                    slide_dict[sent] = KeywordExtract.get_keywords(sent)
                                    list_key.append(text_pair)
                                    all_sents += sent + ". "
                        frame_counter += 1
                if hasattr(shape, "image"):
                  all_sents += shape.alt_text + ". "
                  text_pair = [shape.image, shape.alt_text]
                  slide_dict[shape.alt_text] = KeywordExtract.get_keywords(shape.alt_text)
                  images.append(shape.image)
                  body_text.append(text_pair)
                  list_key.append(text_pair)
            title_dict[title] = list_key
    return body_text, slide_dict, title_dict, all_sents

# # Example usage
# pptx_file = "./my_env/garbage_collection.pptx"
# body_text, slide_dict, title_dict = extract_notes_from_slides(pptx_file)

# print("\nBody Text:")
# print(body_text)
    
# print("\nSlide Dict")
# print(slide_dict)

# print("\nTitle Dict")
# print(title_dict)